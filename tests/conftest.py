"""Shared fixtures and test file generators for doc2md tests."""

from __future__ import annotations

import tempfile
from pathlib import Path

import pymupdf
import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches


@pytest.fixture
def tmp_output_dir():
    """Provide a temporary output directory."""
    with tempfile.TemporaryDirectory(prefix="doc2md_test_") as d:
        yield Path(d)


@pytest.fixture
def sample_pdf(tmp_path: Path) -> Path:
    """Create a minimal PDF test file with text and a table."""
    pdf_path = tmp_path / "test_sample.pdf"
    doc = pymupdf.open()

    # Page 1: text content
    page = doc.new_page()
    text = "Test Document Title\n\nThis is paragraph one with some content.\n\nSecond paragraph here."
    page.insert_text((72, 72), text, fontsize=12)

    # Page 2: more text
    page2 = doc.new_page()
    page2.insert_text((72, 72), "Page Two Content\n\nAnother paragraph on page two.", fontsize=12)

    doc.save(str(pdf_path))
    doc.close()
    return pdf_path


@pytest.fixture
def sample_docx(tmp_path: Path) -> Path:
    """Create a minimal DOCX test file with headings, paragraphs, and a table."""
    docx_path = tmp_path / "test_sample.docx"
    doc = Document()

    doc.add_heading("Test Document", level=1)
    doc.add_paragraph("This is a test paragraph with some content.")
    doc.add_heading("Section Two", level=2)
    doc.add_paragraph("Another paragraph under section two.")

    # Add a table
    table = doc.add_table(rows=3, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Alpha"
    table.cell(1, 1).text = "100"
    table.cell(2, 0).text = "Beta"
    table.cell(2, 1).text = "200"

    # Add a bulleted list
    doc.add_paragraph("Item one", style="List Bullet")
    doc.add_paragraph("Item two", style="List Bullet")

    # Add bold and italic text
    p = doc.add_paragraph()
    p.add_run("Bold text").bold = True
    p.add_run(" and ")
    p.add_run("italic text").italic = True

    doc.save(str(docx_path))
    return docx_path


@pytest.fixture
def sample_pptx(tmp_path: Path) -> Path:
    """Create a minimal PPTX test file with slides, text, and notes."""
    pptx_path = tmp_path / "test_sample.pptx"
    prs = Presentation()

    # Slide 1: Title slide
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Presentation Title"
    slide.placeholders[1].text = "Subtitle text here"

    # Slide 2: Content slide with bullet points
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Slide Two Title"
    body = slide2.placeholders[1]
    body.text = "First bullet point"
    p = body.text_frame.add_paragraph()
    p.text = "Second bullet point"
    p.level = 0

    # Add speaker notes
    notes_slide = slide2.notes_slide
    notes_slide.notes_text_frame.text = "These are speaker notes for slide 2."

    # Slide 3: Slide with a table
    slide_layout = prs.slide_layouts[5]  # Blank
    slide3 = prs.slides.add_slide(slide_layout)
    rows, cols = 3, 2
    table = slide3.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(4), Inches(2)).table
    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"
    table.cell(1, 0).text = "Cell 1"
    table.cell(1, 1).text = "Cell 2"
    table.cell(2, 0).text = "Cell 3"
    table.cell(2, 1).text = "Cell 4"

    prs.save(str(pptx_path))
    return pptx_path
