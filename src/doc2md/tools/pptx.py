"""PPTX to Markdown converter using python-pptx."""

from __future__ import annotations

import logging

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from doc2md.models import ConversionResult
from doc2md.utils.file_handler import (
    cleanup_temp_dir,
    resolve_output_path,
    resolve_source_file,
    write_markdown,
)
from doc2md.utils.markdown import (
    clean_text,
    count_words,
    format_table,
    generate_frontmatter,
)
from doc2md.utils.metadata import build_metadata

logger = logging.getLogger(__name__)


def _extract_text_frame(text_frame) -> list[str]:
    """Extract text from a text frame, preserving bullets."""
    lines = []
    for para in text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        level = para.level
        indent = "  " * level
        # Check if it's a bullet point
        if level > 0:
            lines.append(f"{indent}- {text}")
        else:
            lines.append(text)
    return lines


def _shape_to_markdown(shape, image_counter: list[int]) -> list[str]:
    """Convert a shape to markdown lines."""
    lines = []

    if shape.has_text_frame:
        lines.extend(_extract_text_frame(shape.text_frame))

    if shape.has_table:
        table = shape.table
        headers = [clean_text(cell.text.strip()) for cell in table.rows[0].cells]
        rows = []
        for row in list(table.rows)[1:]:
            rows.append([clean_text(cell.text.strip()) for cell in row.cells])
        md = format_table(headers, rows)
        if md:
            lines.append(md)

    if hasattr(shape, "image"):
        try:
            image_counter[0] += 1
            desc = shape.name or f"image {image_counter[0]}"
            lines.append(f"![{desc}](image_{image_counter[0]}.png)")
        except Exception:
            pass

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            lines.extend(_shape_to_markdown(child, image_counter))

    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
        lines.append(f"*[Chart: {shape.name or 'unnamed chart'}]*")

    return lines


def convert_pptx_to_markdown(
    file_path: str | None = None,
    base64_content: str | None = None,
    file_name: str | None = None,
    output_dir: str | None = None,
    output_file_name: str | None = None,
    overwrite: bool = False,
) -> ConversionResult:
    """Convert a PPTX file to Markdown and write to disk."""
    temp_dir = None
    try:
        source_path, source_name, temp_dir = resolve_source_file(
            file_path, base64_content, file_name
        )
        is_base64 = base64_content is not None

        prs = Presentation(str(source_path))
        slide_count = len(prs.slides)
        warnings: list[str] = []
        total_images = 0
        sections: list[str] = []
        image_counter = [0]

        for slide_num, slide in enumerate(prs.slides, 1):
            # Get slide title
            title = ""
            for shape in slide.shapes:
                if shape.has_text_frame and shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                        idx = shape.placeholder_format.idx
                        if idx == 0:  # Title placeholder
                            title = shape.text_frame.text.strip()
                            break

            if not title:
                # Fallback: use first text shape
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        title = shape.text_frame.text.strip()
                        break

            header = f"## Slide {slide_num}: {title}" if title else f"## Slide {slide_num}"
            slide_lines = [header, ""]

            # Process all shapes
            for shape in slide.shapes:
                # Skip the title shape we already extracted
                if (
                    shape.has_text_frame
                    and hasattr(shape, "placeholder_format")
                    and shape.placeholder_format
                    and shape.placeholder_format.idx == 0
                ):
                    continue

                before = image_counter[0]
                shape_md = _shape_to_markdown(shape, image_counter)
                total_images += image_counter[0] - before

                if shape_md:
                    slide_lines.extend(shape_md)
                    slide_lines.append("")

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_lines.append(f"> **Speaker Notes:** {notes}")
                    slide_lines.append("")

            sections.append("\n".join(slide_lines))

        body = "\n\n".join(sections)
        word_count = count_words(body)

        frontmatter = generate_frontmatter(
            source=source_name,
            fmt="pptx",
            slides=slide_count,
            word_count=word_count,
            warnings=warnings if warnings else None,
        )

        md_content = frontmatter + "\n" + body + "\n"

        output_path = resolve_output_path(
            source_name, source_path, output_dir, output_file_name, overwrite, is_base64
        )
        write_markdown(output_path, md_content)

        metadata = build_metadata(
            source_format="pptx",
            slide_count=slide_count,
            word_count=word_count,
            has_images=total_images > 0,
            image_count=total_images,
            warnings=warnings,
        )

        return ConversionResult(
            success=True,
            output_path=str(output_path),
            file_name=output_path.name,
            source_file=source_name,
            metadata=metadata,
        )

    except Exception as e:
        logger.error("PPTX conversion failed: %s", e)
        return ConversionResult(
            success=False,
            source_file=file_name or (file_path or "unknown"),
            error=str(e),
        )
    finally:
        cleanup_temp_dir(temp_dir)
